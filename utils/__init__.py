"""
Reference: https://github.com/OpenDCAI/DataFlow-Agent
"""
import os
import math
from pathlib import Path
from typing import Dict, List
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from utils.logger import get_logger

log = get_logger(__name__)


def get_project_root() -> Path:
    """
    获取项目根目录路径
    
    Returns:
        Path对象，指向项目根目录
    """
    # 获取当前文件的目录
    current_file = Path(__file__).resolve()
    # 向上查找，直到找到包含特定标记文件的目录（如requirements.txt）
    for parent in [current_file.parent] + list(current_file.parents):
        if (parent / "requirements.txt").exists() or (parent / "pyproject.toml").exists():
            return parent
    # 如果找不到，返回当前文件所在目录的父目录（通常是项目根目录）
    return current_file.parent.parent

def pixels_to_inches(pixels: int, dpi: int = 96) -> float:
    """
    将像素转换为英寸
    
    Args:
        pixels: 像素值
        dpi: 每英寸点数，默认96
        
    Returns:
        英寸值
    """
    return pixels / dpi

def setup_presentation_size(prs, slide_width_px: int = 1024, slide_height_px: int = 1024):
    """设置PPT尺寸"""
    prs.slide_width = Inches(pixels_to_inches(slide_width_px))
    prs.slide_height = Inches(pixels_to_inches(slide_height_px))
    
    return slide_width_px, slide_height_px

def add_text_element(slide, element: Dict):
    """添加文本元素到幻灯片"""
    bbox = element.get('bbox', [0, 0, 100, 50])
    text = element.get('text', '')
    text_level = element.get('text_level')
    
    # 计算位置和大小
    left = pixels_to_inches(bbox[0])
    top = pixels_to_inches(bbox[1])
    width = pixels_to_inches(bbox[2] - bbox[0])
    height = pixels_to_inches(bbox[3] - bbox[1])
    
    # 计算字体大小
    font_size = calculate_font_size(text, bbox, text_level)
    
    log.info(f"添加文本框:")
    log.info(f"  位置: [{bbox[0]}, {bbox[1]}, {bbox[2]}, {bbox[3]}] 像素")
    log.info(f"  英寸坐标: left={left:.2f}, top={top:.2f}, width={width:.2f}, height={height:.2f}")
    log.info(f"  文本内容: {text[:30]}{'...' if len(text) > 30 else ''}")
    log.info(f"  文本级别: {text_level}, 字体大小: {font_size}pt")
    
    # 添加文本框
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width) * 1.2, Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    # 设置文本内容
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    
    # 设置字体样式
    paragraph.font.size = Pt(font_size)
    paragraph.font.name = "Comic Sans MS"
    
    # 根据文本级别设置样式
    if text_level == 1:
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        log.info("  样式: 标题(加粗、居中)")
    elif text_level == 2:
        paragraph.font.bold = True
        log.info("  样式: 子标题(加粗)")
    else:
        log.info("  样式: 正文")
    
    return textbox

def add_image_element(slide, element: Dict):
    """添加图片元素到幻灯片"""
    bbox = element.get('bbox', [0, 0, 100, 100])
    img_path = element.get('img_path', '')
    
    # 计算位置和大小
    left = pixels_to_inches(bbox[0])
    top = pixels_to_inches(bbox[1])
    width = pixels_to_inches(bbox[2] - bbox[0])
    height = pixels_to_inches(bbox[3] - bbox[1])
    
    log.info(f"添加图片:")
    log.info(f"  位置: [{bbox[0]}, {bbox[1]}, {bbox[2]}, {bbox[3]}] 像素")
    log.info(f"  英寸坐标: left={left:.2f}, top={top:.2f}, width={width:.2f}, height={height:.2f}")
    log.info(f"  图片路径: {img_path}")
    log.info(f"  图片尺寸: {bbox[2]-bbox[0]}x{bbox[3]-bbox[1]} 像素")
    
    # 检查图片文件是否存在
    if os.path.exists(img_path):
        try:
            log.info("  图片文件存在，正在添加...")
            result = slide.shapes.add_picture(
                img_path,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            log.info("  图片添加成功")
            return result
        except Exception as e:
            log.error(f"  添加图片时出错: {e}")
            return add_image_placeholder(slide, bbox, f"Error: {str(e)}")
    else:
        log.warning("  图片文件不存在，使用占位符")
        return add_image_placeholder(slide, bbox, "Image not found")

def add_image_placeholder(slide, bbox: List[int], message: str):
    """添加图片占位符"""
    left = pixels_to_inches(bbox[0])
    top = pixels_to_inches(bbox[1])
    width = pixels_to_inches(bbox[2] - bbox[0])
    height = pixels_to_inches(bbox[3] - bbox[1])
    
    # 添加矩形作为占位符
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    # 添加说明文字
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.text = message
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.name = "Comic Sans MS"
    text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return shape

def calculate_font_size(text: str, bbox: List[int], text_level: int = None) -> int:
    """
    根据文本框大小、文字内容和文本级别计算合适的字体大小
    """
    # 计算文本框的宽度和高度（像素）
    width = bbox[2] - bbox[0]
    height = bbox[3] - bbox[1]
    
    # 根据文本级别设置基础字体大小
    if text_level == 1:  # 主标题
        base_size = min(height * 0.8, 44)
    elif text_level == 2:  # 子标题
        base_size = min(height * 0.7, 32)
    else:  # 正文
        base_size = min(height * 0.6, 24)
    
    # 根据文本长度调整
    char_count = len(text)
    if char_count > 0:
        chars_per_line = max(1, width / (base_size * 0.6))
        lines_needed = math.ceil(char_count / chars_per_line)
        
        max_lines = max(1, height / (base_size * 1.1))
        if lines_needed > max_lines:
            base_size = base_size * (max_lines / lines_needed)
    
    # 限制字体大小范围
    font_size = max(8, min(base_size, 72))
    
    return int(font_size)