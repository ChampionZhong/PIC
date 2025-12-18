"""
Figure to ppt service

Reference: https://github.com/OpenDCAI/DataFlow-Agent
"""

import os
import time
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, field
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, pixels_to_inches
from pptx.dml.color import RGBColor

from utils import (
    get_project_root,
    setup_presentation_size,
    add_text_element,
    add_image_element,
    pixels_to_inches as utils_pixels_to_inches,
)
from tools.sam_tool import segment_layout_boxes, free_sam_model
from tools.bg_tool import (
    local_tool_for_raster_to_svg,
    local_tool_for_bg_remove,
    free_bg_rm_model,
)
from tools.mineru_tool import recursive_mineru_layout, svg_to_emf, run_aio_two_step_extract
from tools.image_edit import generate_layout_image_async
from utils.logger import get_logger

log = get_logger(__name__)

# Layout template generation prompt
TEMPLATE_EDIT_PROMPT = (
    "Transform the original image into a pure layout made ONLY of solid colored blocks:\n"
    "1. Keep only the outermost rectangles and arrows (if they exist).\n"
    "2. Delete everything inside them: all titles, subtitles, texts, icons, illustrations, and any inner shapes.\n"
    "3. Turn each remaining outer shape into a solid color block; remove borders if possible.\n"
    "4. Keep the layout exactly the same: same positions, sizes, alignment, and spacing.\n"
    "5. Do NOT add any text, labels, or symbols anywhere.\n"
    "Finally, output a description of this empty color-block template (no text content at all)."
)


@dataclass
class FigureToPPTState:
    """Figure to ppt pipeline state"""
    fig_draft_path: str
    fig_layout_path: str
    output_dir: str = ""
    mineru_port: Optional[int] = None  # Optional, only used for local service mode
    sam_checkpoint: str = ""
    bg_rm_model: str = ""
    mask_detail_level: int = 3
    figure_complex: str = "medium"
    
    # 输出字段
    layout_items: List[Dict[str, Any]] = field(default_factory=list)
    fig_mask: List[Dict[str, Any]] = field(default_factory=list)
    ppt_path: str = ""


def _ensure_result_path(output_dir: str = "") -> str:
    """
    Ensure the output directory for this pipeline
    
    Args:
        output_dir: The specified output directory, if empty, generate automatically
        
    Returns:
        The output directory path
    """
    if output_dir:
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        return str(output_path.resolve())
    
    # Generate output directory automatically
    root = get_project_root()
    ts = int(time.time())
    base_dir = (root / "outputs" / "figure_to_ppt" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    return str(base_dir)


async def process_layout_sam(state: FigureToPPTState) -> FigureToPPTState:
    """
    SAM segment layout elements (empty template image)
    """
    img_path = Path(state.fig_layout_path)
    if not img_path.exists():
        log.error(f"[process_layout_sam] fig_layout_path not found: {img_path}")
        raise FileNotFoundError(f"Layout image not found: {img_path}")
    
    base_dir = Path(_ensure_result_path(state.output_dir))
    out_dir = base_dir / "layout_items"
    out_dir.mkdir(parents=True, exist_ok=True)
    
    # Determine SAM checkpoint path
    if state.sam_checkpoint:
        sam_ckpt = state.sam_checkpoint
    else:
        sam_ckpt = str(get_project_root() / "models" / "SAM" / "sam_b.pt")
    
    log.info(f"[process_layout_sam] Using SAM checkpoint: {sam_ckpt}")
    
    # SAM segment + filter + crop sub-images
    layout_items = segment_layout_boxes(
        image_path=str(img_path),
        output_dir=str(out_dir),
        checkpoint=sam_ckpt,
        min_area=200,
        min_score=0.0,
        iou_threshold=0.2,
        top_k=15,
        nms_by="mask",
    )
    log.info(f"[process_layout_sam] SAM segment result: {len(layout_items)} layout elements")
    
    # Free SAM model
    free_sam_model(checkpoint=sam_ckpt)
    
    # Get layout image size
    try:
        layout_img = Image.open(str(img_path))
        layout_w, layout_h = layout_img.size
    except Exception as e:
        log.error(f"[process_layout_sam] Failed to open layout image: {e}")
        layout_w, layout_h = 1024, 1024
    
    # Each layout PNG to SVG -> EMF, and supplement pixel coordinate bbox_px
    for idx, it in enumerate(layout_items):
        png_path = it.get("png_path")
        if not png_path:
            continue
        
        # Normalize bbox -> pixel bbox
        bbox = it.get("bbox")
        if bbox and len(bbox) == 4:
            x1n, y1n, x2n, y2n = bbox
            x1 = int(round(x1n * layout_w))
            y1 = int(round(y1n * layout_h))
            x2 = int(round(x2n * layout_w))
            y2 = int(round(y2n * layout_h))
            if x2 > x1 and y2 > y1:
                it["bbox_px"] = [x1, y1, x2, y2]
        
        # PNG -> SVG
        svg_path = out_dir / f"layout_{idx}.svg"
        svg_abs = local_tool_for_raster_to_svg({
            "image_path": png_path,
            "output_svg": str(svg_path),
            "colormode": "color",
            "hierarchical": "stacked",
            "mode": "spline",
        })
        it["svg_path"] = svg_abs
        
        # SVG -> EMF
        emf_path = out_dir / f"layout_{idx}.emf"
        try:
            emf_abs = svg_to_emf(svg_abs, str(emf_path))
            it["emf_path"] = emf_abs
        except Exception as e:
            log.error(f"[process_layout_sam] svg_to_emf失败: {svg_abs}, {e}")
            it["emf_path"] = None
    
    state.layout_items = layout_items
    state.output_dir = str(base_dir)  # Update output_dir
    log.info(f"[process_layout_sam] Total {len(layout_items)} layout elements generated")
    return state


async def process_mask_generator(state: FigureToPPTState) -> FigureToPPTState:
    """
    MinerU parse content elements (with content image)
    """
    img_path = Path(state.fig_draft_path)
    if not img_path.exists():
        log.error(f"[process_mask_generator] fig_draft_path not found: {img_path}")
        raise FileNotFoundError(f"Draft image not found: {img_path}")
    
    base_dir = Path(state.output_dir or _ensure_result_path())
    out_dir = base_dir / "mineru_recursive"
    out_dir.mkdir(parents=True, exist_ok=True)
    
    log.info(f"[process_mask_generator] MinerU output directory: {out_dir}")
    log.info(f"[process_mask_generator] mask_detail_level: {state.mask_detail_level}")
    
    # Check MinerU configuration
    from tools.mineru_tool import _get_mineru_api_config
    api_key, api_url, use_api = _get_mineru_api_config()
    if use_api:
        log.info(f"[process_mask_generator] Using MinerU API mode: {api_url}")
    else:
        log.info(f"[process_mask_generator] Using MinerU local service mode, port: {state.mineru_port}")
        if state.mineru_port is None:
            state.mineru_port = 8002  # Default port
    
    # MinerU recursive parse
    mineru_items = await recursive_mineru_layout(
        image_path=str(img_path),
        port=state.mineru_port,
        max_depth=state.mask_detail_level,
        output_dir=out_dir,
    )
    log.info(f"[process_mask_generator] MinerU parse result: {len(mineru_items)} elements")
    
    # Get top-level image size
    top_img = Image.open(state.fig_draft_path)
    top_w, top_h = top_img.size
    
    # Icon original image output directory
    icons_raw_dir = base_dir / "icons_raw"
    icons_raw_dir.mkdir(parents=True, exist_ok=True)
    
    fig_mask = []
    icon_count = 0
    text_count = 0
    
    # Determine details threshold based on figure_complex
    details_map = {"easy": 1, "hard": 10, "medium": 5}
    details = details_map.get(state.figure_complex, 5)
    
    # If the number of elements returned by MinerU is too few, use SAM layout for secondary split
    if len(mineru_items) <= details and state.layout_items:
        log.info(
            f"[process_mask_generator] mineru_items count ({len(mineru_items)}) <= {details}, "
            f"Use SAM layout ({len(state.layout_items)} elements) for secondary MinerU split"
        )
        
        sub_root_dir = base_dir / "mineru_sub_images"
        sub_root_dir.mkdir(parents=True, exist_ok=True)
        
        for layout_idx, layout_it in enumerate(state.layout_items):
            bbox_px = layout_it.get("bbox_px") or layout_it.get("bbox")
            if not bbox_px or len(bbox_px) != 4:
                continue
            
            lx1, ly1, lx2, ly2 = bbox_px
            
            # Boundary crop
            lx1 = max(0, min(top_w, int(round(lx1))))
            ly1 = max(0, min(top_h, int(round(ly1))))
            lx2 = max(0, min(top_w, int(round(lx2))))
            ly2 = max(0, min(top_h, int(round(ly2))))
            if lx2 <= lx1 or ly2 <= ly1:
                continue
            
            # Crop sub-images
            try:
                sub_img = top_img.crop((lx1, ly1, lx2, ly2))
                sub_dir = sub_root_dir / f"layout_{layout_idx}"
                sub_dir.mkdir(parents=True, exist_ok=True)
                sub_path = sub_dir / f"sam_sub_{layout_idx}.png"
                sub_img.save(sub_path)
                
                # Call MinerU for sub-images
                sub_blocks = await run_aio_two_step_extract(str(sub_path), port=state.mineru_port)
                sub_w, sub_h = sub_img.size
                
                # Map to the whole image coordinate system
                for blk_idx, blk in enumerate(sub_blocks):
                    blk_type = blk.get("type", "").lower()
                    bbox_norm = blk.get("bbox")
                    text = blk.get("text") or blk.get("content") or ""
                    
                    if not bbox_norm or len(bbox_norm) != 4:
                        continue
                    
                    sx1n, sy1n, sx2n, sy2n = bbox_norm
                    sx1 = int(round(sx1n * sub_w))
                    sy1 = int(round(sy1n * sub_h))
                    sx2 = int(round(sx2n * sub_w))
                    sy2 = int(round(sy2n * sub_h))
                    
                    gx1 = lx1 + sx1
                    gy1 = ly1 + sy1
                    gx2 = lx1 + sx2
                    gy2 = ly1 + sy2
                    
                    gx1 = max(0, min(top_w, gx1))
                    gy1 = max(0, min(top_h, gy1))
                    gx2 = max(0, min(top_w, gx2))
                    gy2 = max(0, min(top_h, gy2))
                    
                    if gx2 <= gx1 or gy2 <= gy1:
                        continue
                    
                    px_bbox = [gx1, gy1, gx2, gy2]
                    
                    if blk_type in ["title", "text"]:
                        fig_mask.append({
                            "type": "text",
                            "bbox": px_bbox,
                            "text": text,
                            "text_level": 1 if blk_type == "title" else None,
                            "page_idx": 0,
                        })
                        text_count += 1
                    else:
                        crop = top_img.crop((gx1, gy1, gx2, gy2))
                        icon_path = icons_raw_dir / f"blk_sub_{layout_idx}_{blk_idx}.png"
                        crop.save(icon_path)
                        fig_mask.append({
                            "type": "image",
                            "bbox": px_bbox,
                            "img_path": str(icon_path),
                            "page_idx": 0,
                        })
                        icon_count += 1
            except Exception as e:
                log.error(f"[process_mask_generator] Failed to process sub-image layout_idx={layout_idx}: {e}")
                continue
    else:
        # Normal path: directly process MinerU results
        for idx, it in enumerate(mineru_items):
            elem_type = it.get("type", "").lower()
            bbox = it.get("bbox")
            text = (it.get("text") or it.get("content") or "").strip()
            
            if not bbox or len(bbox) != 4:
                continue
            
            x1n, y1n, x2n, y2n = bbox
            x1 = int(round(x1n * top_w))
            y1 = int(round(y1n * top_h))
            x2 = int(round(x2n * top_w))
            y2 = int(round(y2n * top_h))
            
            if x2 <= x1 or y2 <= y1:
                continue
            
            px_bbox = [x1, y1, x2, y2]
            
            if text:
                fig_mask.append({
                    "type": "text",
                    "bbox": px_bbox,
                    "text": text,
                    "text_level": 1 if elem_type == "title" else None,
                    "page_idx": 0,
                })
                text_count += 1
            else:
                crop = top_img.crop((x1, y1, x2, y2))
                icon_path = icons_raw_dir / f"blk_{idx}.png"
                crop.save(icon_path)
                fig_mask.append({
                    "type": "image",
                    "bbox": px_bbox,
                    "img_path": str(icon_path),
                    "page_idx": 0,
                })
                icon_count += 1
    
    state.fig_mask = fig_mask
    log.info(
        f"[process_mask_generator] Total {len(fig_mask)} elements parsed "
        f"(text={text_count}, image={icon_count})"
    )
    return state


async def process_icon_bg_remover(state: FigureToPPTState) -> FigureToPPTState:
    """
    Remove icon background
    """
    base_dir = Path(state.output_dir)
    icons_dir = base_dir / "icons"
    icons_dir.mkdir(parents=True, exist_ok=True)
    
    # Determine RMBG model path
    if state.bg_rm_model:
        bg_model_path = state.bg_rm_model
    else:
        bg_model_path = str(get_project_root() / "models" / "RMBG")
    
    img_cnt = 0
    for item in state.fig_mask:
        if item.get("type") in ["image", "table"]:
            img_cnt += 1
            output_path = local_tool_for_bg_remove({
                "image_path": item.get("img_path"),
                "model_path": bg_model_path,
                "output_dir": str(icons_dir),
            })
            if output_path:
                item["img_path"] = output_path
                log.info(f"[process_icon_bg_remover] Background removed: {output_path}")
            else:
                log.warning(
                    f"[process_icon_bg_remover] Failed to remove background: {item.get('img_path')}"
                )
    
    log.info(f"[process_icon_bg_remover] Processed {img_cnt} image/table elements")
    
    # Free RMBG model
    try:
        free_bg_rm_model(model_path=bg_model_path)
        log.info("[process_icon_bg_remover] RMBG-2.0 model freed")
    except Exception as e:
        log.error(f"[process_icon_bg_remover] Failed to free model: {e}")
    
    return state


async def process_ppt_generation(state: FigureToPPTState) -> FigureToPPTState:
    """
    Generate PPT
    """
    output_dir = Path(state.output_dir or _ensure_result_path())
    output_dir.mkdir(parents=True, exist_ok=True)
    
    timestamp = int(time.time())
    ppt_filename = f"presentation_{timestamp}.pptx"
    ppt_path = output_dir / ppt_filename
    
    prs = Presentation()
    
    # Set PPT size
    img = Image.open(state.fig_draft_path)
    width_px, height_px = img.size
    slide_width_px, slide_height_px = setup_presentation_size(prs, width_px, height_px)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    def _add_layout_emf(slide, item) -> bool:
        """Add EMF from layout_item to slide by pixel bbox"""
        emf_path = item.get("emf_path")
        if not emf_path or not os.path.exists(emf_path):
            if emf_path:
                log.warning(f"[process_ppt_generation] emf_path not found: {emf_path}")
            return False
        
        bbox = item.get("bbox_px") or item.get("bbox")
        if not bbox or len(bbox) != 4:
            log.warning(f"[process_ppt_generation] layout_item missing valid bbox: {item}")
            return False
        
        x1, y1, x2, y2 = bbox
        if x2 <= x1 or y2 <= y1:
            log.warning(f"[process_ppt_generation] Invalid pixel bbox coordinates: {bbox}")
            return False
        
        # Pixel -> inches
        left_in = utils_pixels_to_inches(x1)
        top_in = utils_pixels_to_inches(y1)
        width_in = utils_pixels_to_inches(x2 - x1)
        height_in = utils_pixels_to_inches(y2 - y1)
        
        try:
            slide.shapes.add_picture(
                emf_path,
                Inches(left_in),
                Inches(top_in),
                Inches(width_in),
                Inches(height_in),
            )
            return True
        except Exception as e:
            log.error(f"[process_ppt_generation] Failed to add_picture EMF: {emf_path}, {e}")
            return False
    
    # Page 1: Complete combined page
    slide_main = prs.slides.add_slide(blank_slide_layout)
    background = slide_main.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Render layout_items (background layer)
    layout_drawn = 0
    for item in state.layout_items:
        if _add_layout_emf(slide_main, item):
            layout_drawn += 1
    
    # Render fig_mask (content layer)
    img_drawn = 0
    text_drawn = 0
    for element in state.fig_mask:
        elem_type = element.get("type", "")
        if elem_type == "text":
            add_text_element(slide_main, element)
            text_drawn += 1
        elif elem_type in ["image", "table"]:
            add_image_element(slide_main, element)
            img_drawn += 1
    
    # Page 2: Only EMF debug page
    slide_emf = prs.slides.add_slide(blank_slide_layout)
    bg2 = slide_emf.background
    fill2 = bg2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 255, 255)
    for item in state.layout_items:
        _add_layout_emf(slide_emf, item)
    
    # Page 3: Complete original image page
    slide_full = prs.slides.add_slide(blank_slide_layout)
    bg3 = slide_full.background
    fill3 = bg3.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(255, 255, 255)
    
    left_in = utils_pixels_to_inches(0)
    top_in = utils_pixels_to_inches(0)
    width_in = utils_pixels_to_inches(width_px)
    height_in = utils_pixels_to_inches(height_px)
    
    slide_full.shapes.add_picture(
        state.fig_draft_path,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    
    # Save PPT
    prs.save(str(ppt_path))
    state.ppt_path = str(ppt_path)
    
    log.info(f"[process_ppt_generation] PPT generated successfully: {ppt_path}")
    log.info(f"[process_ppt_generation] Layout elements: {len(state.layout_items)}, drawn: {layout_drawn}")
    log.info(f"[process_ppt_generation] Content elements: {len(state.fig_mask)}, text={text_drawn}, image={img_drawn}")
    
    return state


async def generate_layout_image(
    fig_draft_path: str,
    output_dir: str,
    api_url: str = None,
    api_key: str = None,
    model: str = None,
) -> str:
    """
    Automatically generate layout template image
    
    Args:
        fig_draft_path: Path to the image with content
        output_dir: Output directory
        api_url: API address (optional, from environment variables)
        api_key: API Key (optional, from environment variables)
        model: Model name (optional, from environment variables)
        
    Returns:
        Path to the generated layout image
    """
    import time
    
    # Get API configuration from environment variables
    if not api_url:
        api_url = os.getenv("GEMINI_API_URL", "http://123.129.219.111:3000/v1")
    if not api_key:
        api_key = os.getenv("GEMINI_API_KEY", os.getenv("DF_API_KEY", ""))
    if not model:
        model = os.getenv("GEMINI_MODEL", "gemini-3-pro-image-preview")
    
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    timestamp = int(time.time())
    layout_filename = f"layout_{timestamp}.png"
    layout_save_path = str(output_path / layout_filename)
    
    log.info(f"[generate_layout_image] Generating layout template from: {fig_draft_path}")
    log.info(f"  API URL: {api_url}")
    log.info(f"  Model: {model}")
    log.info(f"  Output: {layout_save_path}")
    
    try:
        await generate_layout_image_async(
            prompt=TEMPLATE_EDIT_PROMPT,
            image_path=fig_draft_path,
            save_path=layout_save_path,
            api_url=api_url,
            api_key=api_key,
            model=model,
            aspect_ratio="16:9",
            resolution="2K",
            timeout=300,
        )
        log.info(f"[generate_layout_image] Layout template generated: {layout_save_path}")
        return layout_save_path
    except Exception as e:
        log.error(f"[generate_layout_image] Failed to generate layout template: {e}", exc_info=True)
        raise


async def run_figure_to_ppt_pipeline_with_auto_layout(
    fig_draft_path: str,
    output_dir: str = "",
    mineru_port: Optional[int] = None,
    sam_checkpoint: str = "",
    bg_rm_model: str = "",
    mask_detail_level: int = 3,
    figure_complex: str = "medium",
    api_url: str = None,
    api_key: str = None,
    model: str = None,
) -> str:
    """
    Main pipeline function (version with auto layout image generation)
    
    Args:
        fig_draft_path: Path to the image with content
        output_dir: Output directory
        mineru_port: MinerU service port
        sam_checkpoint: SAM model checkpoint path (optional)
        bg_rm_model: RMBG model path (optional)
        mask_detail_level: MinerU recursive depth
        figure_complex: Image complexity ("easy"/"medium"/"hard")
        api_url: API address (optional)
        api_key: API Key (optional)
        model: Model name (optional)
        
    Returns:
        Path to the generated PPT file
    """
    log.info("[run_figure_to_ppt_pipeline_with_auto_layout] Starting pipeline with auto layout generation")
    
    # Ensure output directory exists
    output_path = Path(output_dir or _ensure_result_path())
    output_path.mkdir(parents=True, exist_ok=True)
    
    # 1. Automatically generate layout image
    log.info("[run_figure_to_ppt_pipeline_with_auto_layout] Step 1: Generating layout template")
    fig_layout_path = await generate_layout_image(
        fig_draft_path=fig_draft_path,
        output_dir=str(output_path),
        api_url=api_url,
        api_key=api_key,
        model=model,
    )
    
    # 2. Run the original pipeline
    state = await run_figure_to_ppt_pipeline(
        fig_draft_path=fig_draft_path,
        fig_layout_path=fig_layout_path,
        output_dir=str(output_path),
        mineru_port=mineru_port,
        sam_checkpoint=sam_checkpoint,
        bg_rm_model=bg_rm_model,
        mask_detail_level=mask_detail_level,
        figure_complex=figure_complex,
    )
    
    return state.ppt_path


async def run_figure_to_ppt_pipeline(
    fig_draft_path: str,
    fig_layout_path: str,
    output_dir: str = "",
    mineru_port: Optional[int] = None,
    sam_checkpoint: str = "",
    bg_rm_model: str = "",
    mask_detail_level: int = 3,
    figure_complex: str = "medium",
) -> FigureToPPTState:
    """
    Main pipeline function
    
    Args:
        fig_draft_path: Path to the image with content
        fig_layout_path: Path to the empty template image
        output_dir: Output directory (optional)
        mineru_port: MinerU service port
        sam_checkpoint: SAM model checkpoint path (optional)
        bg_rm_model: RMBG model path (optional)
        mask_detail_level: MinerU recursive depth
        figure_complex: Image complexity ("easy"/"medium"/"hard")
        
    Returns:
        FigureToPPTState object, containing the generated PPT path and information
    """
    log.info("[run_figure_to_ppt_pipeline] Starting pipeline")
    
    # Create state object
    state = FigureToPPTState(
        fig_draft_path=fig_draft_path,
        fig_layout_path=fig_layout_path,
        output_dir=output_dir,
        mineru_port=mineru_port,
        sam_checkpoint=sam_checkpoint,
        bg_rm_model=bg_rm_model,
        mask_detail_level=mask_detail_level,
        figure_complex=figure_complex,
    )
    
    try:
        # Execute pipeline steps
        state = await process_layout_sam(state)
        state = await process_mask_generator(state)
        state = await process_icon_bg_remover(state)
        state = await process_ppt_generation(state)
        
        log.info(f"[run_figure_to_ppt_pipeline] Pipeline executed successfully, PPT path: {state.ppt_path}")
        return state
        
    except Exception as e:
        log.error(f"[run_figure_to_ppt_pipeline] Pipeline execution failed: {e}", exc_info=True)
        raise

